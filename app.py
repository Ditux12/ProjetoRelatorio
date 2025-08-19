import argparse
from collections import OrderedDict
import pandas as pd
import matplotlib.pyplot as plt
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import unicodedata
from io import BytesIO
import numpy as np
from datetime import datetime, timedelta
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
import tempfile
import os
import base64
from fastapi.responses import JSONResponse

OPINION_CATEGORIES = ["Artigo de Opinião", "Comentário"]
IGNORE_CATEGORIES = ["Desporto"]
ARTOPINION_CATEGORIES = ["Artigo de Opinião", "Comentário"]

TITLE_FONT = 32
SUBTITLE_FONT = 18
BODY_FONT = 12

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ICON_PATH = os.path.join(BASE_DIR, "static", "u4.png")
IMAGE_PATH = os.path.join(BASE_DIR, "static", "u23.png")

app= FastAPI()

@app.post("/generate-report")
async def generate_report(file: UploadFile = File(...)):
    try:
        # Guardar Excel temporariamente
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp_in:
            content = await file.read()
            tmp_in.write(content)
            tmp_in_path = tmp_in.name

        # Criar ficheiro temporário para o PPTX
        tmp_out = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        tmp_out_path = tmp_out.name
        tmp_out.close()

        # Usar a tua função principal
        main(tmp_in_path, tmp_out_path)

        # Ler o PPTX gerado em bytes
        with open(tmp_out_path, "rb") as f:
            pptx_bytes = f.read()

        # Converter para Base64
        pptx_b64 = base64.b64encode(pptx_bytes).decode("utf-8")

        # Apagar ficheiros temporários
        os.remove(tmp_in_path)
        os.remove(tmp_out_path)

        # Retornar como JSON
        return JSONResponse(content={"file_base64": pptx_b64})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/ping")
async def ping():
    return JSONResponse(content={"status": "awake"})



def set_cell_border(cell, color=RGBColor(0,0,0), width=12700):
    """Define bordas para uma célula usando XML (width em EMUs, 12700 ≈ 0.127mm)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ('lnL', 'lnR', 'lnT', 'lnB'):
        ln = tcPr.find(qn(f'a:{border_name}'))
        if ln is None:
            ln = parse_xml(f'<a:{border_name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
            tcPr.append(ln)

        # Define largura
        ln.set('w', str(width))

        # Remove preenchimento antigo
        solidFill = ln.find(qn('a:solidFill'))
        if solidFill is not None:
            ln.remove(solidFill)

        # Cor
        solidFill = parse_xml(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{color[0]:02X}{color[1]:02X}{color[2]:02X}"/>'

            f'</a:solidFill>'
        )
        ln.append(solidFill)

def read_excel(path):
    import openpyxl
    # Ler normal para DataFrame
    df = pd.read_excel(path)
    df = df.rename(columns=lambda c: c.strip())
    if 'Data de publicação' in df.columns:
        df['Data de publicação'] = pd.to_datetime(df['Data de publicação']).dt.date

    # Abrir com openpyxl para pegar links
    wb = openpyxl.load_workbook(path, data_only=True)
    ws = wb.active

    # Encontrar índice da coluna "Título"
    titulo_col_idx = None
    for idx, cell in enumerate(ws[1], start=1):
        if str(cell.value).strip().lower() == "título":
            titulo_col_idx = idx
            break

    link_map = {}
    if titulo_col_idx:
        for row_idx in range(2, ws.max_row + 1):
            cell = ws.cell(row=row_idx, column=titulo_col_idx)
            if cell.hyperlink:
                link_map[row_idx - 2] = cell.hyperlink.target  # índice no DataFrame

    # Criar coluna Link
    df["Link"] = df.index.map(link_map).fillna("")
    return df




def set_slide_background(slide, rgb_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)

def add_icon_to_slide(slide, icon_path):
    slide.shapes.add_picture(icon_path, Inches(0.2), Inches(0.2), height=Inches(0.9))

def add_image_to_slide(slide, image_path):
    left = Inches(-0.69)
    top = Inches(1.52)
    width = Inches(10.69)
    height = Inches(5.98)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)

from pptx.util import Pt

def normalize(text):
    """Remove acentos, espaços extras e converte para minúsculas."""
    return ''.join(c for c in unicodedata.normalize('NFD', text)
                   if unicodedata.category(c) != 'Mn').strip().lower()

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def add_index_slide(prs, sections, slide_refs, page_numbers):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Índice"
    
    # Formatar título
    title_tf = slide.shapes.title.text_frame
    p = title_tf.paragraphs[0].font
    p.name = 'Barlow'
    p.size = Pt(TITLE_FONT)
    p.color.rgb = RGBColor(255, 255, 255)

    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(0.5)

    for idx, sec in enumerate(sections):
        y = top + Inches(0.6 * idx)
        shape = slide.shapes.add_textbox(left, y, width, height)
        tf = shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]

        # Pegar título e número da página
        title = sec['title']
        page_num = page_numbers.get(title, "?")  # se não encontrar, usa '?'

        # Criar linha estilo: Título..............Número
        dots = '.' * (60 - len(title))  # ajusta a quantidade de pontos
        para.text = f"{title}{dots}{page_num}"

# Configuração do estilo do índice
        para.font.name = 'Barlow'
        para.font.size = Pt(20)
        para.font.color.rgb = RGBColor(255, 255, 255)  # azul
        para.font.underline = True  # sublinhado


        # Adicionar hyperlink interno se existir
        if title in slide_refs:
            shape.click_action.target_slide = slide_refs[title]

    return slide


def _make_leader_line(label: str, page_num: int, width: int = 70) -> str:
    s_page = str(page_num)
    dots = max(2, width - len(label) - len(s_page))
    return f"{label}{'.' * dots}{s_page}"




def create_pie_chart(df):
    counts = df['Meio'].value_counts()
    labels = counts.index.tolist()
    sizes = counts.values.tolist()
    colors = plt.cm.tab20.colors[:len(labels)]

    fig, ax = plt.subplots(figsize=(5,5))
    fig.patch.set_facecolor('#404040')
    ax.set_facecolor('#404040')

    wedges, texts, autotexts = ax.pie(
        sizes,
        labels=labels,
        autopct='%1.1f%%',
        colors=colors,
        startangle=90,
        wedgeprops={'edgecolor':'white', 'linewidth':1},
        textprops={'color':'white', 'fontsize':10},
        pctdistance=0.6,   # percentagens dentro do gráfico
        labeldistance=1.05 # labels fora
    )

    ax.axis('equal')

    # Colocar percentagens pequenas verticalmente
    for i, autotext in enumerate(autotexts):
        if sizes[i] / sum(sizes) < 0.05:  # fatias pequenas <5%
            autotext.set_rotation(90)  # vertical
            autotext.set_verticalalignment('center')
            autotext.set_horizontalalignment('center')

    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True, dpi=300)
    plt.close(fig)
    buf.seek(0)
    return buf



def build_overview_table(prs, stats, pie_img_bytes):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Overview"
    title_tf = slide.shapes.title.text_frame
    p = title_tf.paragraphs[0].font
    p.name = 'Barlow'
    p.size = Pt(TITLE_FONT)
    p.color.rgb = RGBColor(255, 255, 255)

    # --- Tabela ---
    rows = len(stats['by_category'])+1
    cols = 3
    table = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.5), Inches(4.5), Inches(3)
    ).table
    headers = ["Categoria", "Nº Notícias", "Circulação"]

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(64, 64, 64)

    for i, (cat, vals) in enumerate(stats['by_category'].items(), start=1):
        table.cell(i, 0).text = cat
        table.cell(i, 1).text = str(vals['count'])
        table.cell(i, 2).text = f"{vals['circ']:,}"
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(64, 64, 64)

    # --- Gráfico ---
    slide.shapes.add_picture(
        pie_img_bytes,
        Inches(5.5),   # X
        Inches(1.29),  # Y
        width=Inches(3.94),
        height=Inches(3.8)
    )

    # --- Caixa de texto ---
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(5), Inches(1)).text_frame
    tx.text = (
        f"—Total de notícias: {stats['total_rows']}\n"
        f"—Circulação total: {stats['total_circ']:,}\n"
        f"—AAV total: {stats['total_aav']:,}"
    )

    for p in tx.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)

    tx.paragraphs[0].font.name = 'Barlow'
    tx.paragraphs[0].font.size = Pt(SUBTITLE_FONT)
    tx.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    return slide

def add_slide_numbers(prs):
    for i, slide in enumerate(prs.slides, start=1):
        left = prs.slide_width - Inches(1) - Inches(0.2)
        top = prs.slide_height - Inches(0.3) - Inches(0.2)
        width = Inches(1)
        height = Inches(0.3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(i)
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.RIGHT

from pptx.enum.dml import MSO_FILL_TYPE

def add_table_slide(prs, category_name, items, rows_per_slide=6):
    total_count = len(items)
    total_circ = int(items['Circulação'].sum()) if total_count > 0 else 0

    # Slide de introdução do tema
    slide_intro = prs.slides.add_slide(prs.slide_layouts[5])
    slide_intro.shapes.title.text = f"{category_name}"
    title_tf = slide_intro.shapes.title.text_frame
    p = title_tf.paragraphs[0].font
    p.name = 'Barlow'
    p.size = Pt(TITLE_FONT)
    p.color.rgb = RGBColor(255, 255, 255)

    # Caixa 1 - Total de notícias
    shape1 = slide_intro.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(1.78), Inches(3), Inches(1.5))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(64, 64, 64)
    shape1.line.color.rgb = RGBColor(64, 64, 64)
    shape1.line.width = Pt(1)
    tf1 = shape1.text_frame
    tf1.clear()
    p1 = tf1.add_paragraph()
    p1.text = f"Total de notícias: {total_count}"
    p1.font.size = Pt(20)
    p1.font.name = 'Impact'
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER

    # Caixa 2 - Circulação acumulada
    shape2 = slide_intro.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), Inches(4.53), Inches(3), Inches(1.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(64, 64, 64)
    shape2.line.color.rgb = RGBColor(64, 64, 64)
    shape2.line.width = Pt(1)
    tf2 = shape2.text_frame
    tf2.clear()
    p2 = tf2.add_paragraph()
    p2.text = f"Circulação acumulada: {total_circ:,}"
    p2.font.size = Pt(20)
    p2.font.name = 'Impact'
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER

    # Linha vertical no meio
    line = slide_intro.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.14), Inches(0), Inches(0.03), Inches(7.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 255, 255)
    line.line.fill.background()  # sem borda

    # Ajuste do título à direita
    title_shape = slide_intro.shapes.title
    title_shape.left = Inches(0.3)
    title_shape.top = Inches(3.28)
    title_shape.width = Inches(3)
    title_shape.height = Inches(0.71)
    title_shape.text = f"{category_name}"
    title_tf = title_shape.text_frame
    p = title_tf.paragraphs[0].font
    p.name = 'Barlow'
    p.size = Pt(25)
    p.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 🔹 Definir colunas extras só para Artigos de opinião e Comentários
    if category_name in ["Artigos de opinião", "Comentários"]:
        extra_cols = ["Autor", "Instituição"]
    else:
        extra_cols = []
    base_cols = ["Meio", "Data de publicação", "Título", "Publicação", "Circulação"] + extra_cols

    # Função auxiliar para criar a tabela em cada chunk
    def _create_table_for_chunk(chunk, slide_title):
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Ajuste do título do slide
        title_shape = slide.shapes.title
        title_shape.left = Inches(1.77)
        title_shape.top = Inches(0.29)
        title_shape.width = Inches(7.05)
        title_shape.height = Inches(0.71)
        title_shape.text = slide_title
        title_tf = title_shape.text_frame
        p = title_tf.paragraphs[0].font
        p.name = 'Barlow'
        p.size = Pt(25)
        p.color.rgb = RGBColor(255, 255, 255)
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # --- TABELA ---
        rows = len(chunk) + 1
        cols = len(base_cols)
        table_height = Inches(1.5) if rows <= 4 else Inches(5)

        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), table_height).table

        # Cabeçalhos
        for j, col_name in enumerate(base_cols):
            cell = table.cell(0, j)
            cell.text = col_name
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            set_cell_border(cell, RGBColor(0, 0, 0))
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Dados da tabela
        for i, (_, r) in enumerate(chunk.iterrows(), start=1):
            vals = [
                r.get("Meio", ""),
                r.get("Data de publicação", "").isoformat()
                if hasattr(r.get("Data de publicação", ""), "isoformat")
                else "",
                r.get("Título", ""),
                r.get("Publicação", ""),
                str(int(r.get("Circulação", 0))) if pd.notna(r.get("Circulação", 0)) else "0",
            ]
            # 🔹 Se categoria for artigos/comentários, adiciona Autor e Instituição
            if "Autor" in base_cols:
                vals.append(r.get("Autor", ""))
            if "Instituição" in base_cols:
                vals.append(r.get("Instituição", ""))

            for j, val in enumerate(vals):
                cell = table.cell(i, j)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.word_wrap = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                set_cell_border(cell, RGBColor(0, 0, 0))

                if base_cols[j] == "Título" and pd.notna(r.get("Link")) and r.get("Link") != "":
                    run = p.add_run()
                    run.text = str(val)
                    run.hyperlink.address = r.get("Link")
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.underline = True
                else:
                    p.text = str(val)

        return slide

    # Criar os slides de dados (com Tema Secundário em todas as categorias)
    if "Tema Secundário" in items.columns and items["Tema Secundário"].notna().any():
        grouped = items.groupby("Tema Secundário")
        for tema, subset in grouped:
            if pd.isna(tema) or str(tema).strip() == "":
                continue  # 🔹 evita criar slides "Categoria — nan"
            subset = subset.sort_values('Data de publicação', ascending=False)
            for start in range(0, len(subset), rows_per_slide):
                chunk = subset.iloc[start:start + rows_per_slide]
                _create_table_for_chunk(chunk, f"{category_name} — {tema}")
    else:
        items = items.sort_values('Data de publicação', ascending=False)
        for start in range(0, len(items), rows_per_slide):
            chunk = items.iloc[start:start + rows_per_slide]
            _create_table_for_chunk(chunk, category_name)

    return slide_intro











def _render_table_slide(prs, category_name, rows, base_cols):
    """Função auxiliar para desenhar uma tabela no slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{category_name}"

    table = slide.shapes.add_table(
        len(rows) + 1, len(base_cols), Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    ).table

    # Cabeçalhos
    for j, col_name in enumerate(base_cols):
        cell = table.cell(0, j)
        cell.text = col_name
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Linhas
    for i, row in enumerate(rows, start=1):
        vals = [
            row.get("Meio", ""),
            row.get("Data de publicação", "").isoformat() if hasattr(row.get("Data de publicação", ""), "isoformat") else str(row.get("Data de publicação", "")),
            row.get("Título", ""),
            row.get("Publicação", ""),
            str(int(row.get("Circulação", 0))) if pd.notna(row.get("Circulação", 0)) else "0",
        ]
        extra_cols = [c for c in base_cols if c not in ["Meio", "Data de publicação", "Título", "Publicação", "Circulação"]]
        for c in extra_cols:
            vals.append(row.get(c, ""))

        for j, val in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.text_frame.word_wrap = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)


def add_cover_slide(prs, title, icon_path, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, (64, 64, 64))
    add_icon_to_slide(slide, icon_path)
    add_image_to_slide(slide, image_path)
    title_shape = slide.shapes.title
    title_shape.text = title
    title_tf = title_shape.text_frame
    p = title_tf.paragraphs[0].font
    p.name = 'Barlow'
    p.size = Pt(48)
    p.color.rgb = RGBColor(255, 255, 255)
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_shape.left = Inches(1)
    title_shape.top = Inches(0.5)
    title_shape.width = Inches(8)
    title_shape.height = Inches(1.5)
    
    # ===== Adicionar caixa de texto com intervalo de datas =====
    from datetime import datetime, timedelta

    hoje = datetime.today()
    sete_dias_antes = hoje - timedelta(days=7)

    # Formatar no estilo: 31 de julho – 1 de agosto
    meses = ["janeiro", "fevereiro", "março", "abril", "maio", "junho",
             "julho", "agosto", "setembro", "outubro", "novembro", "dezembro"]

    inicio_str = f"{sete_dias_antes.day} de {meses[sete_dias_antes.month-1]}"
    fim_str = f"{hoje.day} de {meses[hoje.month-1]}"

    date_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5))
    tf = date_box.text_frame
    tf.clear()
    p_date = tf.add_paragraph()
    p_date.text = f"{inicio_str} – {fim_str}"
    p_date.font.size = Pt(16)
    p_date.font.name = 'Calibri'
    p_date.font.color.rgb = RGBColor(255, 255, 255)
    p_date.alignment = PP_ALIGN.CENTER
    # ========================================================

    return slide


def add_closing_slide(prs, icon_path, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide, (64, 64, 64))
    add_icon_to_slide(slide, icon_path)
    add_image_to_slide(slide, image_path)
    title_shape = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.5))
    tf = title_shape.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Fim do Relatório"
    p.font.name = 'Barlow'
    p.font.size = Pt(48)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    return slide

def main(input_path, output_path):
    df = read_excel(input_path)
    expected_cols = ['Meio','Data de publicação','Título','Publicação','Circulação',
                     'Tema Principal','Tema Secundário','Autor','Instituição','AAV']
    for c in expected_cols:
        if c not in df.columns:
            df[c] = None

    # Ajustar categorias
    df['Categoria_final'] = df['Tema Principal'].replace({
        "Artigo de Opinião": "Artigos de opinião",
        "Comentário": "Comentários"
    })

    df.loc[~df['Tema Principal'].isin(OPINION_CATEGORIES), 'Categoria_final'] = \
        df.loc[~df['Tema Principal'].isin(OPINION_CATEGORIES), 'Tema Principal']

    df = df[~df['Categoria_final'].isin(IGNORE_CATEGORIES)]

    # Ordenar categorias: normais primeiro, opinião/comentário sempre por último
    opinion_cats = ["Artigos de opinião", "Comentários"]
    all_categories = df['Categoria_final'].dropna().unique().tolist()
    normal_cats = [cat for cat in all_categories if cat not in opinion_cats]
    CATEGORY_ORDER = sorted(normal_cats) + [cat for cat in opinion_cats if cat in all_categories]

    total_rows = len(df)
    total_circ = int(df['Circulação'].sum()) if total_rows > 0 else 0
    total_aav = int(df['AAV'].sum()) if 'AAV' in df.columns and df['AAV'].notna().any() else 0

    by_category = OrderedDict()
    for cat in CATEGORY_ORDER:
        sub = df[df['Categoria_final'] == cat]
        if len(sub) > 0:
            by_category[cat] = {'count': len(sub), 'circ': int(sub['Circulação'].sum())}

    stats = {
        'total_rows': total_rows,
        'total_circ': total_circ,
        'total_aav': total_aav,
        'by_category': by_category
    }

    prs = Presentation()

    # 1. Slide de capa
    add_cover_slide(prs, "Relatório de notícias semanal", ICON_PATH, IMAGE_PATH)

    # 2. Overview
    pie_buf = create_pie_chart(df)
    overview_slide = build_overview_table(prs, stats, pie_buf)
    
    # Criar slides de categorias e guardar referência da INTRODUÇÃO
    slide_refs = {"Overview": overview_slide}
    for cat in CATEGORY_ORDER:
        if cat not in by_category:
            continue
        items = df[df['Categoria_final'] == cat].copy()
        items = items.sort_values('Data de publicação', ascending=False)

        if cat == "Artigos de opinião e comentários":
            first_slide = add_table_slide(prs, cat, items, extra_cols=["Autor", "Instituição"])
        else:
            first_slide = add_table_slide(prs, cat, items)

        slide_refs[cat] = first_slide

    # Calcular números de página antes de mexer na ordem dos slides
    page_numbers = {cat: list(prs.slides).index(slide) + 2 
                    for cat, slide in slide_refs.items() if slide is not None}

    # Criar índice (links para intro e overview)
    index_slide = add_index_slide(
        prs,
        [{'title': k} for k in slide_refs.keys()],
        slide_refs,
        page_numbers
    )

    # Mover índice para depois da capa
    prs.slides._sldIdLst.insert(1, prs.slides._sldIdLst[-1])

    # Slide final
    add_closing_slide(prs, ICON_PATH, IMAGE_PATH)

    # Fundo + ícone
    for slide in prs.slides:
        set_slide_background(slide, (64, 64, 64))
        add_icon_to_slide(slide, ICON_PATH)

    # Números no canto inferior
    add_slide_numbers(prs)

    prs.save(output_path)


if __name__ == "__main__":
    input_path = "C:/Users/Diogo/Desktop/conversor/excel121.xlsx"
    output_path = "Relatorio_Tabelas.pptx"
    main(input_path, output_path)
    print(f"PPTX gerado: {output_path}")


